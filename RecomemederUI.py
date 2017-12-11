from pptx import Presentation
from tkinter import *
from PIL import Image, ImageTk
import urllib.request
import dict_mem as dict_mem
import check_similarity as similiarity
from sklearn.metrics.pairwise import cosine_similarity
from sklearn.feature_extraction.text import TfidfVectorizer
import pandas as pd
import re, math
from collections import Counter

WORD = re.compile(r'\w+')

def get_cosine(vec1, vec2):
     intersection = set(vec1.keys()) & set(vec2.keys())
     numerator = sum([vec1[x] * vec2[x] for x in intersection])

     sum1 = sum([vec1[x]**2 for x in vec1.keys()])
     sum2 = sum([vec2[x]**2 for x in vec2.keys()])
     denominator = math.sqrt(sum1) * math.sqrt(sum2)

     if not denominator:
        return 0.0
     else:
        return float(numerator) / denominator

def text_to_vector(text):
     words = WORD.findall(text)
     return Counter(words)


# text_runs will be populated with a list of strings,
# one for each text run in presentation
def get_ppt_text(filename):
    prs = Presentation(filename)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return text_runs

class InputPage(Frame):
    def __init__(self, parent, controller,*args, **kwargs):
        Frame.__init__(self,parent,*args,**kwargs)
        #set frame's app
        self.app = controller

        #create upload presentation frame
        presentationFrame = Frame(self)
        presentationLabel = Label(presentationFrame, text="Enter your presentation filename:",font="Arial 14 bold")
        presentationLabel.pack()
        self.presentationEntry = Entry(presentationFrame, width = 20)
        self.presentationEntry.insert(END, '202 Final Demo')

        self.presentationEntry.pack()
        presentationFrame.pack(pady=10)

        #create sentiment frame
        sentimentsFrame = Frame(self)
        sentimentsLabel = Label(sentimentsFrame, text="Select the intended sentiment(s):",font="Arial 14 bold")
        sentimentsLabel.pack()
        sentimentsFrame.pack(pady=10)

        # create dictionary with sentiments
        self.sentiments = {'Joy':"",
        'Sadness':"",
        'Anger':"",
        'Disgust':"",
        'Surprise':"",
        'Anticipation':"",
        'Trust':"",
        'Sarcasm':"",
        'Shame':"",
        'Fear':"",
        'Triumphant':"",
        'Confused':""}

        #create checkboxes for sentiments
        for item in self.sentiments:
            value = IntVar()
            checkbox = Checkbutton(sentimentsFrame, text=item, variable=value)
            checkbox.var = value
            self.sentiments[item]=checkbox
            checkbox.pack()

        #create audience frame
        audienceFrame = Frame(self)
        audienceLabel = Label(audienceFrame, text="Select the intended audience(s):", font="Arial 14 bold")
        audienceLabel.pack()
        audienceFrame.pack(pady=10)

        # create dictionary with audience
        self.audience = {'Friends':"",
        'Class-Peers':"",
        'Class-Acquaintances-And-Faculty':"",
        'Work-Internal':"",
        'Work-Clients':""}

        #create checkboxes for audience
        for item in self.audience:
            value = IntVar()
            checkbox = Checkbutton(audienceFrame, text=item, variable=value)
            checkbox.var = value
            self.audience[item]=checkbox
            checkbox.pack()

        submitButton = Button(self, text="Reco-meme!", command=lambda : self.sendRequest())
        submitButton.pack()

    # send request for selected inputs and presentation
    def sendRequest(self):
        selectedSentiments = []
        selectedAudience = []

        presentationFilename = self.presentationEntry.get() + '.pptx'
        print("presentation filename: ", presentationFilename)
        presentationText = get_ppt_text(presentationFilename)
        print("presentation text: ", presentationText)

        for item in self.sentiments:
            value = self.sentiments[item].var.get()
            if value != 0:
                selectedSentiments.append(item)
        print("selected sentiments: ", selectedSentiments)

        for item in self.audience:
            value = self.audience[item].var.get()
            if value != 0:
                selectedAudience.append(item)
        print("selected audience: ", selectedAudience)

        sentiment_audience_list = similiarity.check_similarity(selectedSentiments, selectedAudience)

        final_list = []
        if len(sentiment_audience_list) <= 3:
            final_list
            memeData = pd.read_csv("dataset.csv")
            for element in sentiment_audience_list:
                final_list.append((memeData['URL'].iloc[int(element)], memeData['Subject Matter'].iloc[int(element)]))

        else:
            final_list = self.getcosinesimilarity(presentationText, self.get_description(sentiment_audience_list))


        #create the output page TODO: call this function somewhere with inputs on the recommended memes to create the output
        self.createOutputPage(final_list)

        return presentationFilename, selectedSentiments, selectedAudience

    # TODO: add parameter for recommended memes information

    def getcosinesimilarity(self, ppt_text_list, url_caption_subject_list):

        simi_url_list = []
        print ('Calculating cosine similarity between captions and presentation text')
        ppt_text = " ".join(str(x) for x in ppt_text_list)
        for url_caption in url_caption_subject_list:
            url, caption, subject = url_caption

            vector1 = text_to_vector(ppt_text)
            vector2 = text_to_vector(caption)
            cosine = get_cosine(vector1, vector2)

            simi_url_list.append((cosine, url, subject))
        final = sorted(simi_url_list, key=lambda x: x[0], reverse=True)
        new_list = []
        for element in final[:3]:
            cos, url, subject = element
            new_list.append((url, subject))
        return new_list

    def get_description(self, list_of_index):
        desc = []
        memeData = pd.read_csv("dataset.csv")
        for i in list_of_index:
            desc.append((memeData['URL'].iloc[int(i)], memeData['Caption'].iloc[int(i)], memeData['Subject Matter'].iloc[int(i)]))
        return desc

    def createOutputPage(self, urls):
        outputFrame = OutputPage(self.app.container, self.app, urls)
        outputFrame.pack()
        outputFrame.tkraise()

        #destroy the input frame
        self.destroy()

class OutputPage(Frame):
    def __init__(self, parent, controller, urls, *args, **kwargs):
        Frame.__init__(self,parent,*args,**kwargs)

        #set frame's app
        self.app = controller

        #create recommendations frame
        recommendationsFrame = Frame(self)
        recommendationsLabel = Label(recommendationsFrame, text="Your Top " + str(len(urls)) + " Recommended Meme(s)",font="Arial 16 bold")
        recommendationsLabel.pack(fill=X)
        recommendationsFrame.pack(fill=X)

        #create dictionary for recommended memes, TODO: need recommended memes and somehow get their urls and name

        #retrieve the images
        for i in range(len(urls)):
            memeFrame = Frame(self)
            memeFrame.pack(fill=X, pady=10)

            url, subject = urls[i]

            #download meme image. TODO: need to replace with meme's url and meme's filename
            urllib.request.urlretrieve(url,"testimagename" + str(i) + ".jpg")

            print ('URL', url)

            basewidth = 150
            image = Image.open("testimagename" + str(i) + ".jpg") #TODO: input meme's filename
            wpercent = (basewidth / float(image.size[0]))
            hsize = int((float(image.size[1]) * float(wpercent)))
            image = image.resize((basewidth, hsize), Image.ANTIALIAS)
            photo = ImageTk.PhotoImage(image)

            memeImage = Label(memeFrame, image=photo)
            memeImage.image = photo
            memeImage.pack(fill=X)

            memeLabel = Label(memeFrame, text=subject) #TODO: replace with meme's metadata
            memeLabel_url = Label(memeFrame, text=url, fg="blue", cursor="hand2")

            #memeLabel_url = Label(memeFrame, text=url) #TODO: replace with meme's metadata
            memeLabel.pack(fill=X)
            memeLabel_url.pack(fill=X)

        #new recommendation button
        newButton = Button(self, text="Reco-meme Again!", command=lambda : self.createInputPage())
        newButton.pack()

    def createInputPage(self):
        inputFrame = InputPage(self.app.container, self.app)
        inputFrame.pack()
        inputFrame.tkraise()

        #destroy the output frame
        self.destroy()

class RecomemederUI(Tk):
    def __init__(self, *args, **kwargs):
        Tk.__init__(self, *args, **kwargs)

        self.title("The Reco-MEME-der")
        self.minsize(width=800, height=800)
        self.maxsize(width=1000, height=1000)
        self.configure(background='#E1ECF4')

        #create the main container to that other frames will be raised into
        self.container = Frame(self, background='white')
        self.container.pack(side='top', fill='both', expand=True)

        #create the input page
        inputFrame = InputPage(self.container, self)
        inputFrame.pack()

if __name__ == '__main__':
    # create inverted lists for audience and sentiments
    try:
        print("Loading, please wait ....")
        print("Creating a beautiful UI, just for you ....")
        dict_mem.controller()

        # Retrieve user choices
        app = RecomemederUI()
        app.mainloop()
    except Exception:
        print('Error')