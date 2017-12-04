from pptx import Presentation
from tkinter import *
prs = Presentation('test.pptx')

# text_runs will be populated with a list of strings,
# one for each text run in presentation
def get_ppt_text():
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)

    print text_runs

def main():
    root = Tk()
    root.title("The Reco-MEME-der")
    root.minsize(width=800, height=800)
    root.maxsize(width=1000, height=1000)
    root.configure(background='#E1ECF4')

    # Add a grid
    mainframe = Frame(root)
    mainframe.grid(column=0, row=0, sticky=(N, W, E, S))
    mainframe.columnconfigure(0, weight=1)
    mainframe.rowconfigure(0, weight=1)
    mainframe.pack(pady=100, padx=100)

    # Create a Tkinter variable
    tkvar = StringVar(root)

    # Dictionary with options
    choices = {'Joy',
'Sadness',
'Anger',
'Disgust',
'Surprise',
'Anticipation',
'Trust',
'Sarcasm',
'Shame',
'Fear',
'Triumphant',
'Confused'}
    tkvar.set('Happy')  # set the default option

    popupMenu = OptionMenu(mainframe, tkvar, *choices)
    Label(mainframe, text="Choose a sentiment").grid(row=1, column=1)
    popupMenu.grid(row=2, column=1)

    # on change dropdown value
    def change_dropdown(*args):
        print(tkvar.get())

    # link function to change dropdown
    tkvar.trace('w', change_dropdown)

    # Add a grid
    mainframe = Frame(root)
    mainframe.grid(column=0, row=0, sticky=(N, W, E, S))
    mainframe.columnconfigure(0, weight=1)
    mainframe.rowconfigure(0, weight=1)
    mainframe.pack(pady=100, padx=100)

    # Create a Tkinter variable
    tkvar = StringVar(root)

    # Dictionary with options
    choices = {'Friends', 'Class - Peers', 'Class - Internals and Faculty', 'Work - Internals', 'Work - Clients'}
    tkvar.set('Happy')  # set the default option

    popupMenu = OptionMenu(mainframe, tkvar, *choices)
    Label(mainframe, text="Choose a sentiment").grid(row=1, column=1)
    popupMenu.grid(row=2, column=1)

    # on change dropdown value
    def change_dropdown(*args):
        print(tkvar.get())

    # link function to change dropdown
    tkvar.trace('w', change_dropdown)

    button = Button(root, bg='gray', text="Reco-Meme!", width=40, height=3)
    button.pack()

    root.mainloop()

if __name__=="__main__":
    main()

